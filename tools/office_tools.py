"""Native Microsoft Office file tools for Kara.

Paths are constrained by the same read/write roots and sensitive-file policy as
Kara's general file tools. Office files are manipulated directly through OOXML
libraries; Microsoft Office does not need to be open.
"""
from __future__ import annotations

import csv
import json
from datetime import date, datetime
from io import StringIO
from pathlib import Path
from typing import Any

import config
from tools.file_tools import _resolve_path


def _json(payload: dict[str, Any]) -> str:
    return json.dumps(payload, indent=2, ensure_ascii=False)


def _error(message: str) -> str:
    return _json({"ok": False, "error": message})


def _new_output_path(path: str, extension: str, overwrite: bool) -> Path:
    target = _resolve_path(path, config.FILE_WRITE_ROOTS, purpose="write")
    if target.suffix.casefold() != extension:
        raise ValueError(f"Output path must end in {extension}.")
    if target.exists() and not overwrite:
        raise FileExistsError(f"File already exists: {target}")
    if target.exists() and target.is_dir():
        raise IsADirectoryError(f"Output path is a directory: {target}")
    target.parent.mkdir(parents=True, exist_ok=True)
    return target


def create_word_document(
    path: str,
    title: str,
    content: str,
    overwrite: bool = False,
) -> str:
    """Create a real Word .docx document from a title and plain-text paragraphs.

    Args:
        path: Destination .docx path inside an allowed write root.
        title: Optional document title rendered with Word's Title style.
        content: Paragraph text; blank lines separate paragraphs.
        overwrite: Replace an existing file only when explicitly true.

    Returns:
        JSON confirmation with paragraph count and saved file size.
    """
    try:
        from docx import Document

        target = _new_output_path(path, ".docx", overwrite)
        document = Document()
        if title.strip():
            document.add_heading(title.strip(), level=0)
        paragraphs = [block.strip() for block in content.split("\n\n") if block.strip()]
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
        document.save(target)
        # Structural verification: reopen the OOXML package through python-docx.
        verified = Document(target)
        return _json(
            {
                "ok": True,
                "path": str(target),
                "type": "word",
                "paragraph_count": len(verified.paragraphs),
                "size_bytes": target.stat().st_size,
            }
        )
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not create Word document: {exc}")


def append_word_text(path: str, content: str) -> str:
    """Append plain-text paragraphs to an existing Word .docx document.

    Args:
        path: Existing .docx path inside an allowed write root.
        content: Paragraph text; blank lines separate paragraphs.

    Returns:
        JSON confirmation after reopening and verifying the saved document.
    """
    try:
        from docx import Document

        target = _resolve_path(path, config.FILE_WRITE_ROOTS, purpose="write")
        if target.suffix.casefold() != ".docx":
            raise ValueError("Word document path must end in .docx.")
        if not target.exists() or not target.is_file():
            raise ValueError(f"Word document does not exist: {target}")
        document = Document(target)
        additions = [block.strip() for block in content.split("\n\n") if block.strip()]
        if not additions:
            raise ValueError("content cannot be empty.")
        for paragraph in additions:
            document.add_paragraph(paragraph)
        document.save(target)
        verified = Document(target)
        return _json(
            {
                "ok": True,
                "path": str(target),
                "type": "word",
                "paragraphs_added": len(additions),
                "paragraph_count": len(verified.paragraphs),
                "size_bytes": target.stat().st_size,
            }
        )
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not append Word text: {exc}")


def _spreadsheet_value(value: str) -> Any:
    stripped = value.strip()
    if stripped == "":
        return None
    if stripped.casefold() in {"true", "false"}:
        return stripped.casefold() == "true"
    try:
        return int(stripped)
    except ValueError:
        try:
            return float(stripped)
        except ValueError:
            return value


def _json_cell(value: Any) -> Any:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return str(value)


def create_excel_workbook(
    path: str,
    data_csv: str,
    sheet_name: str = "Sheet1",
    overwrite: bool = False,
) -> str:
    """Create a real Excel .xlsx workbook from CSV-formatted text.

    Args:
        path: Destination .xlsx path inside an allowed write root.
        data_csv: RFC-style CSV text. Integers, decimals, and booleans become typed cells.
        sheet_name: Worksheet name (1-31 characters; Excel-invalid characters rejected).
        overwrite: Replace an existing file only when explicitly true.

    Returns:
        JSON confirmation with row/column counts and saved file size.
    """
    try:
        from openpyxl import Workbook, load_workbook
        from openpyxl.styles import Font
        from openpyxl.utils import get_column_letter

        target = _new_output_path(path, ".xlsx", overwrite)
        name = sheet_name.strip()
        if not name or len(name) > 31 or any(char in name for char in "[]:*?/\\"):
            raise ValueError("sheet_name must be 1-31 characters and exclude []:*?/\\.")
        rows = list(csv.reader(StringIO(data_csv)))
        if not rows:
            raise ValueError("data_csv must contain at least one row.")
        if len(rows) > 10_000 or max((len(row) for row in rows), default=0) > 200:
            raise ValueError("Workbook input is limited to 10000 rows and 200 columns.")
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = name
        for row in rows:
            sheet.append([_spreadsheet_value(value) for value in row])
        for cell in sheet[1]:
            cell.font = Font(bold=True)
        sheet.freeze_panes = "A2"
        if sheet.max_row and sheet.max_column:
            sheet.auto_filter.ref = sheet.dimensions
        for index, column_cells in enumerate(sheet.columns, 1):
            width = min(60, max(10, max(len(str(cell.value or "")) for cell in column_cells) + 2))
            sheet.column_dimensions[get_column_letter(index)].width = width
        workbook.save(target)
        verified = load_workbook(target, read_only=True, data_only=False)
        verified_sheet = verified[name]
        result = {
            "ok": True,
            "path": str(target),
            "type": "excel",
            "sheet": name,
            "rows": verified_sheet.max_row,
            "columns": verified_sheet.max_column,
            "size_bytes": target.stat().st_size,
        }
        verified.close()
        return _json(result)
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not create Excel workbook: {exc}")


def set_excel_cell(
    path: str,
    cell: str,
    value: str,
    sheet_name: str = "",
) -> str:
    """Set one cell in an existing Excel .xlsx workbook.

    Args:
        path: Existing .xlsx path inside an allowed write root.
        cell: A1-style cell address such as B2.
        value: New value; integers, decimals, and booleans become typed cells.
        sheet_name: Worksheet to edit; blank selects the active sheet.

    Returns:
        JSON confirmation containing the previous and new values.
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.utils.cell import coordinate_to_tuple

        target = _resolve_path(path, config.FILE_WRITE_ROOTS, purpose="write")
        if target.suffix.casefold() != ".xlsx":
            raise ValueError("Excel workbook path must end in .xlsx.")
        if not target.exists() or not target.is_file():
            raise ValueError(f"Excel workbook does not exist: {target}")
        coordinate = cell.strip().upper()
        row_number, column_number = coordinate_to_tuple(coordinate)
        if row_number > 1_048_576 or column_number > 16_384:
            raise ValueError("Cell address is outside Excel worksheet limits.")
        workbook = load_workbook(target)
        if sheet_name.strip():
            if sheet_name not in workbook.sheetnames:
                workbook.close()
                raise ValueError(f"Worksheet does not exist: {sheet_name}")
            sheet = workbook[sheet_name]
        else:
            sheet = workbook.active
        previous = _json_cell(sheet[coordinate].value)
        typed_value = _spreadsheet_value(value)
        sheet[coordinate] = typed_value
        selected_sheet = sheet.title
        workbook.save(target)
        workbook.close()
        verified = load_workbook(target, read_only=True, data_only=False)
        saved_value = _json_cell(verified[selected_sheet][coordinate].value)
        verified.close()
        return _json(
            {
                "ok": True,
                "path": str(target),
                "type": "excel",
                "sheet": selected_sheet,
                "cell": coordinate,
                "previous": previous,
                "value": saved_value,
                "size_bytes": target.stat().st_size,
            }
        )
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not set Excel cell: {exc}")


def _parse_slides_json(slides_json: str) -> list[dict[str, Any]]:
    try:
        data = json.loads(slides_json)
    except json.JSONDecodeError as exc:
        raise ValueError(f"slides_json is not valid JSON: {exc}") from exc
    if not isinstance(data, list) or not data:
        raise ValueError("slides_json must be a non-empty JSON array.")
    if len(data) > 50:
        raise ValueError("A presentation is limited to 50 content slides.")
    slides: list[dict[str, Any]] = []
    for index, item in enumerate(data, 1):
        if not isinstance(item, dict):
            raise ValueError(f"Slide {index} must be a JSON object.")
        title = str(item.get("title") or "").strip()
        bullets = item.get("bullets") or []
        if not title:
            raise ValueError(f"Slide {index} requires a title.")
        if not isinstance(bullets, list) or len(bullets) > 20:
            raise ValueError(f"Slide {index} bullets must be a JSON array of at most 20 items.")
        slides.append({"title": title, "bullets": [str(value) for value in bullets]})
    return slides


def _add_powerpoint_content_slide(presentation: Any, title: str, bullets: list[str]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFB)
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), presentation.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0x00, 0xA6, 0xB2)
    accent.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = title
    title_run.font.name = "Aptos Display"
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x16, 0x22, 0x36)
    body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.65), Inches(11.2), Inches(5.1))
    body_frame = body_box.text_frame
    body_frame.clear()
    body_frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(0x2B, 0x36, 0x47)
        paragraph.space_after = Pt(12)
        paragraph.alignment = PP_ALIGN.LEFT


def create_powerpoint(
    path: str,
    title: str,
    slides_json: str,
    overwrite: bool = False,
) -> str:
    """Create a styled PowerPoint .pptx from a JSON array of titled bullet slides.

    Args:
        path: Destination .pptx path inside an allowed write root.
        title: Presentation title shown on a dedicated title slide.
        slides_json: JSON array like [{"title":"Status","bullets":["Ready"]}].
        overwrite: Replace an existing file only when explicitly true.

    Returns:
        JSON confirmation after reopening and verifying the presentation.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        target = _new_output_path(path, ".pptx", overwrite)
        presentation_title = title.strip()
        if not presentation_title:
            raise ValueError("title cannot be empty.")
        slides = _parse_slides_json(slides_json)
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        cover = presentation.slides.add_slide(presentation.slide_layouts[6])
        cover_fill = cover.background.fill
        cover_fill.solid()
        cover_fill.fore_color.rgb = RGBColor(0x16, 0x22, 0x36)
        title_box = cover.shapes.add_textbox(Inches(1.0), Inches(2.25), Inches(11.3), Inches(1.5))
        frame = title_box.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = presentation_title
        run.font.name = "Aptos Display"
        run.font.size = Pt(42)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        subtitle = cover.shapes.add_textbox(Inches(1.05), Inches(4.05), Inches(6), Inches(0.5))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = "Prepared by Kara"
        subtitle_frame.paragraphs[0].font.name = "Aptos"
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0x35, 0xD0, 0xBA)
        for slide_spec in slides:
            _add_powerpoint_content_slide(
                presentation, slide_spec["title"], slide_spec["bullets"]
            )
        presentation.save(target)
        verified = Presentation(target)
        return _json(
            {
                "ok": True,
                "path": str(target),
                "type": "powerpoint",
                "slide_count": len(verified.slides),
                "size_bytes": target.stat().st_size,
            }
        )
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not create PowerPoint: {exc}")


def append_powerpoint_slide(path: str, title: str, bullets_json: str) -> str:
    """Append one styled bullet slide to an existing PowerPoint .pptx.

    Args:
        path: Existing .pptx path inside an allowed write root.
        title: Slide title.
        bullets_json: JSON string array such as ["First point", "Second point"].

    Returns:
        JSON confirmation after reopening and verifying the presentation.
    """
    try:
        from pptx import Presentation

        target = _resolve_path(path, config.FILE_WRITE_ROOTS, purpose="write")
        if target.suffix.casefold() != ".pptx":
            raise ValueError("PowerPoint path must end in .pptx.")
        if not target.exists() or not target.is_file():
            raise ValueError(f"PowerPoint does not exist: {target}")
        slide_title = title.strip()
        if not slide_title:
            raise ValueError("title cannot be empty.")
        try:
            bullets = json.loads(bullets_json)
        except json.JSONDecodeError as exc:
            raise ValueError(f"bullets_json is not valid JSON: {exc}") from exc
        if not isinstance(bullets, list) or len(bullets) > 20:
            raise ValueError("bullets_json must be a JSON array of at most 20 items.")
        presentation = Presentation(target)
        _add_powerpoint_content_slide(presentation, slide_title, [str(item) for item in bullets])
        presentation.save(target)
        verified = Presentation(target)
        return _json(
            {
                "ok": True,
                "path": str(target),
                "type": "powerpoint",
                "slide_count": len(verified.slides),
                "size_bytes": target.stat().st_size,
            }
        )
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not append PowerPoint slide: {exc}")


def read_office_file(path: str, max_items: int = 200) -> str:
    """Extract bounded structured content from a Word, Excel, or PowerPoint OOXML file.

    Args:
        path: Existing .docx, .xlsx, or .pptx file inside an allowed read root.
        max_items: Maximum paragraphs, rows, or text items returned (1-1000).

    Returns:
        Structured JSON content without launching Microsoft Office.
    """
    try:
        target = _resolve_path(path, config.FILE_READ_ROOTS, purpose="read")
        if not target.exists() or not target.is_file():
            raise ValueError(f"Office file does not exist: {target}")
        limit = max(1, min(int(max_items), 1000))
        extension = target.suffix.casefold()
        if extension == ".docx":
            from docx import Document

            document = Document(target)
            paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
            tables = [
                [[cell.text for cell in row.cells] for row in table.rows]
                for table in document.tables
            ]
            return _json(
                {
                    "ok": True,
                    "type": "word",
                    "path": str(target),
                    "paragraphs": paragraphs[:limit],
                    "tables": tables[:limit],
                    "truncated": len(paragraphs) > limit or len(tables) > limit,
                }
            )
        if extension == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(target, read_only=True, data_only=False)
            sheets: list[dict[str, Any]] = []
            remaining = limit
            truncated = False
            for sheet in workbook.worksheets:
                rows: list[list[Any]] = []
                for row in sheet.iter_rows(values_only=True):
                    if remaining <= 0:
                        truncated = True
                        break
                    rows.append([_json_cell(value) for value in row])
                    remaining -= 1
                sheets.append(
                    {
                        "name": sheet.title,
                        "max_row": sheet.max_row,
                        "max_column": sheet.max_column,
                        "rows": rows,
                    }
                )
                if remaining <= 0:
                    if sheet != workbook.worksheets[-1] or sheet.max_row > len(rows):
                        truncated = True
                    break
            workbook.close()
            return _json(
                {
                    "ok": True,
                    "type": "excel",
                    "path": str(target),
                    "sheets": sheets,
                    "truncated": truncated,
                }
            )
        if extension == ".pptx":
            from pptx import Presentation

            presentation = Presentation(target)
            slides: list[dict[str, Any]] = []
            remaining = limit
            truncated = False
            for number, slide in enumerate(presentation.slides, 1):
                text_items: list[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        candidates = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    else:
                        text = str(getattr(shape, "text", "") or "").strip()
                        candidates = [text] if text else []
                    for text in candidates:
                        if remaining <= 0:
                            truncated = True
                            break
                        text_items.append(text)
                        remaining -= 1
                    if remaining <= 0:
                        break
                slides.append({"number": number, "text": text_items})
                if remaining <= 0:
                    if number < len(presentation.slides):
                        truncated = True
                    break
            return _json(
                {
                    "ok": True,
                    "type": "powerpoint",
                    "path": str(target),
                    "slides": slides,
                    "slide_count": len(presentation.slides),
                    "truncated": truncated,
                }
            )
        raise ValueError("Supported Office extensions are .docx, .xlsx, and .pptx.")
    except (ImportError, ValueError, PermissionError, OSError) as exc:
        return _error(f"Could not read Office file: {exc}")

# --- Registry declaration ------------------------------------------------------
# Consumed by tools.registry; this is the single source of truth for which
# functions in this module are exposed to the model and which of them are safe
# for unattended scheduled runs.
TOOL_GROUP = "office"

TOOLS = [
    read_office_file,
    create_word_document,
    append_word_text,
    create_excel_workbook,
    set_excel_cell,
    create_powerpoint,
    append_powerpoint_slide,
]

SCHEDULED_SAFE = {
    "read_office_file",
}
